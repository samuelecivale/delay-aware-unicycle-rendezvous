#!/usr/bin/env python3
"""
generate_presentation.py

Builds presentation/delay_aware_unicycle_rendezvous.pptx: a 16:9 academic/
engineering slide deck on delay-aware rendezvous of nonholonomic unicycle
multi-robot systems, sourced from:
  - report/delay_aware_unicycle_rendezvous_updated.tex
  - main.m
  - unicycle_only_outputs/figures/*.png
  - unicycle_only_outputs/tables/*.csv

Run:
    source .venv_presentation/bin/activate   # (from repo root)
    python presentation/generate_presentation.py

All repository assets are located via paths relative to this script, so the
script can be re-run from any working directory.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from eq_render import render_eq, render_eq_lines

# ----------------------------------------------------------------------------
# Paths
# ----------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
# NOTE: unicycle_only_outputs/figures|tables|videos/ (no subfolder) are stale leftovers
# from a run predating the paper/vector_field split (confirmed by diffing against
# unicycle_only_outputs/paper/tables/*.csv: numbers match the vector_field controller,
# not 'paper'). The deck's default/described controller is 'paper' (force-projection,
# CONTROLLER_TYPE='paper' in main.m), so all figures/tables/videos are sourced from the
# paper/ subfolder. vector_field/ is used only for the explicit controller-comparison slide.
FIG = os.path.join(REPO, "unicycle_only_outputs", "paper", "figures")
VID = os.path.join(REPO, "unicycle_only_outputs", "paper", "videos")
FIG_VF = os.path.join(REPO, "unicycle_only_outputs", "vector_field", "figures")
VID_VF = os.path.join(REPO, "unicycle_only_outputs", "vector_field", "videos")
ASSETS = os.path.join(HERE, "generated_assets")
OUT_PPTX = os.path.join(HERE, "delay_aware_unicycle_rendezvous.pptx")

os.makedirs(ASSETS, exist_ok=True)

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
BG = RGBColor(0x0A, 0x12, 0x22)          # deep navy background
PANEL = RGBColor(0x11, 0x1D, 0x33)       # slightly lighter panel
PANEL2 = RGBColor(0x16, 0x24, 0x3E)
LINE = RGBColor(0x2A, 0x3B, 0x5C)
WHITE = RGBColor(0xF3, 0xF6, 0xFB)
MUTED = RGBColor(0xA7, 0xB3, 0xC8)
ACCENT_BLUE = RGBColor(0x4F, 0x9C, 0xF9)     # matches matplotlib "full-state" blue
ACCENT_ORANGE = RGBColor(0xFF, 0x9E, 0x4A)   # matches matplotlib "neighbour-only" orange
ACCENT_GREEN = RGBColor(0x4C, 0xD9, 0x88)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_GOLD = RGBColor(0xE8, 0xC5, 0x4A)

FONT = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

SLIDE_NUM = {"n": 0}
SLIDE_TITLES = []  # for outline export


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    SLIDE_NUM["n"] += 1
    return s


def add_rect(slide, left, top, width, height, fill=PANEL, line=LINE, line_w=0.75,
             shadow=False, rounded=False, radius=0.06):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
             line_spacing=1.0, wrap=True, shrink=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, size=18, color=WHITE,
                 bullet_color=ACCENT_BLUE, gap=0.14, font=FONT, bold_lead=None,
                 line_spacing=1.08):
    """items: list of str, or (str, dict) for per-item overrides."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        opts = {}
        if isinstance(item, tuple):
            item, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap * 72)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(opts.get("size", size))
        r.font.color.rgb = opts.get("color", color)
        r.font.bold = opts.get("bold", False)
        r.font.name = font
    return tb


def add_picture_contain(slide, path, left, top, max_w, max_h, align="center"):
    """Place an image inside a bounding box, preserving aspect ratio (letterboxed)."""
    with Image.open(path) as im:
        iw, ih = im.size
    ratio_img = iw / ih
    ratio_box = max_w / max_h
    if ratio_img > ratio_box:
        w = max_w
        h = int(max_w / ratio_img)
    else:
        h = max_h
        w = int(max_h * ratio_img)
    if align == "center":
        l = left + (max_w - w) // 2
    else:
        l = left
    t = top + (max_h - h) // 2
    pic = slide.shapes.add_picture(path, l, t, width=w, height=h)
    return pic, (l, t, w, h)


def add_title(slide, title, kicker=None, y=Inches(0.42)):
    if kicker:
        add_text(slide, Inches(0.55), y, Inches(11.5), Inches(0.32), kicker.upper(),
                  size=13, color=ACCENT_ORANGE, bold=True, font=FONT)
        y = y + Inches(0.32)
    add_text(slide, Inches(0.55), y, Inches(12.2), Inches(0.75), title,
              size=30, color=WHITE, bold=True, font=FONT)
    bar_y = y + Inches(0.68)
    add_rect(slide, Inches(0.55), bar_y, Inches(1.1), Pt(3.2), fill=ACCENT_BLUE, line=None)
    return bar_y + Inches(0.22)


def add_footer(slide, text):
    n = SLIDE_NUM["n"]
    add_text(slide, Inches(0.55), Inches(7.14), Inches(9.0), Inches(0.3), text,
              size=10.5, color=MUTED, font=FONT)
    add_text(slide, Inches(12.3), Inches(7.14), Inches(0.5), Inches(0.3), str(n),
              size=10.5, color=MUTED, font=FONT, align=PP_ALIGN.RIGHT)


def add_arrow_right(slide, left, top, width, height, fill=ACCENT_BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.55
        sh.adjustments[1] = 0.55
    except Exception:
        pass
    return sh


def add_arrow_down(slide, left, top, width, height, fill=ACCENT_BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.55
        sh.adjustments[1] = 0.55
    except Exception:
        pass
    return sh


def add_pill(slide, left, top, width, height, text, fill=PANEL2, text_color=WHITE,
             border=ACCENT_BLUE, size=14, bold=True):
    sh = add_rect(slide, left, top, width, height, fill=fill, line=border, line_w=1.25,
                   rounded=True, radius=0.22)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = FONT
    return sh


def add_eq_picture(slide, path, left, top, max_w, max_h, align="center"):
    return add_picture_contain(slide, path, left, top, max_w, max_h, align=align)


def add_video(slide, movie_path, poster_path, left, top, max_w, max_h, label=None,
              frame_color=ACCENT_BLUE):
    """Insert the trajectory poster frame as a plain static picture (no embedded
    movie shape). Google Slides does not reliably import/play a movie shape
    embedded in a .pptx (add_movie survives round-trips to PowerPoint/Keynote
    but is dropped or frozen on import into Slides), so the deck must stay
    Slides-safe by relying on the poster image alone; the source .mp4 is
    shipped alongside the deck for anyone opening it in PowerPoint/Keynote.
    Returns (picture_shape, (l, t, w, h))."""
    with Image.open(poster_path) as im:
        iw, ih = im.size
    ratio_img = iw / ih
    ratio_box = max_w / max_h
    if ratio_img > ratio_box:
        w = max_w
        h = int(max_w / ratio_img)
    else:
        h = max_h
        w = int(max_h * ratio_img)
    l = left + (max_w - w) // 2
    t = top + (max_h - h) // 2
    pad = Emu(int(0.045 * 914400))
    add_rect(slide, l - pad, t - pad, w + 2 * pad, h + 2 * pad, fill=RGBColor(0xFF, 0xFF, 0xFF),
             line=frame_color, line_w=1.5, rounded=False)
    pic = slide.shapes.add_picture(poster_path, l, t, w, h)
    if label:
        # align with sibling image captions, which sit at (top + max_h): ignore the
        # letterboxed fitted height so captions across a row stay on one baseline.
        add_text(slide, left, top + max_h, max_w, Inches(0.28), label,
                  size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    return pic, (l - pad, t - pad, w + 2 * pad, h + 2 * pad)


def speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


NOTES = {}  # slide index (1-based) -> notes text, exported also to slide_notes.md


# ----------------------------------------------------------------------------
# Equation assets (rendered once via matplotlib mathtext)
# ----------------------------------------------------------------------------
EQ = {}
EQ["unicycle"] = render_eq(
    r"$\dot{x}_i=v_i\cos\theta_i,\qquad \dot{y}_i=v_i\sin\theta_i,\qquad \dot{\theta}_i=\omega_i$",
    "eq_unicycle.png", fontsize=30)
EQ["laplacian"] = render_eq(
    r"$L = D - A$", "eq_laplacian.png", fontsize=32)
EQ["eig_order"] = render_eq(
    r"$0=\lambda_1\leq\lambda_2\leq\cdots\leq\lambda_N=\lambda_{\max}(L)$",
    "eq_eig_order.png", fontsize=24)
EQ["paper_def"] = render_eq(
    r"$u_i\in\mathbb{R}^2$ desired consensus velocity $\quad$ $t_i=(\cos\theta_i,\sin\theta_i)$ heading direction",
    "eq_paper_def.png", fontsize=20)
EQ["paper_law"] = render_eq(
    r"$v_i = k_v\,u_i^{\top}t_i, \qquad \omega_i = k_\theta\,e_\theta$",
    "eq_paper_law.png", fontsize=30)
EQ["e_theta"] = render_eq(
    r"$e_\theta = \mathrm{atan2}(u_{i,y},u_{i,x}) - \theta_i$",
    "eq_e_theta.png", fontsize=20)
EQ["fullstate"] = render_eq(
    r"$\dot p(t) = -k\,L\,p(t-\tau)$", "eq_fullstate.png", fontsize=28)
EQ["neighboronly"] = render_eq(
    r"$u_i(t) = -k\sum_{j\in\mathcal{N}_i} a_{ij}\left(p_i(t)-p_j(t-\tau)\right)$",
    "eq_neighboronly.png", fontsize=24)
EQ["neighboronly_stacked"] = render_eq(
    r"$\dot p(t) = -k\left(D\,p(t) - A\,p(t-\tau)\right)$",
    "eq_neighboronly_stacked.png", fontsize=22)

# derivation steps (slide 6)
EQ["step1a"] = render_eq(
    r"$\dot p(t)=-k\,L\,p(t-\tau)$", "eq_step1a.png", fontsize=26)
EQ["step1b"] = render_eq(
    r"$\dot z_i(t) = -k\lambda_i\,z_i(t-\tau) \qquad (z=V^{\top}p)$",
    "eq_step1b.png", fontsize=26)
EQ["step2"] = render_eq(
    r"$z(t)=e^{st}: \qquad s = -a\,e^{-s\tau}, \qquad a=k\lambda_i$",
    "eq_step2.png", fontsize=24)
EQ["step3"] = render_eq(
    r"$s=j\omega: \qquad j\omega = -a\,e^{-j\omega\tau} = -a(\cos\omega\tau - j\sin\omega\tau)$",
    "eq_step3.png", fontsize=24)
EQ["step4"] = render_eq(
    r"$\cos\omega\tau = 0 \;\Rightarrow\; \omega\tau=\frac{\pi}{2}, \qquad \omega = a\sin\omega\tau = a$",
    "eq_step4.png", fontsize=24)
EQ["step5"] = render_eq(
    r"$\tau_i = \frac{\pi}{2\,k\,\lambda_i} \qquad\Rightarrow\qquad \mathrm{worst\ mode:}\ \lambda_i=\lambda_{\max}(L)$",
    "eq_step5.png", fontsize=24)
EQ["taucrit_big"] = render_eq(
    r"$\tau_{\mathrm{crit}} \;=\; \frac{\pi}{2\,k\,\lambda_{\max}(L)}$",
    "eq_taucrit_big.png", fontsize=44)
EQ["taucrit_small"] = render_eq(
    r"$\tau_{\mathrm{crit}} = \frac{\pi}{2k\lambda_{\max}(L)}$",
    "eq_taucrit_small.png", fontsize=20)
EQ["taucrit_ring"] = render_eq(
    r"$\tau_{\mathrm{crit}} = \frac{\pi}{2(1)(4)} = \frac{\pi}{8} \approx 0.3927\ \mathrm{s}$",
    "eq_taucrit_ring.png", fontsize=22)
EQ["topology_prop"] = render_eq(
    r"$\tau_{\mathrm{crit}} \;\propto\; \frac{1}{\lambda_{\max}(L)}$",
    "eq_topology_prop.png", fontsize=30)
EQ["gain_prop"] = render_eq(
    r"$\tau_{\mathrm{crit}} \;\propto\; \frac{1}{k}$",
    "eq_gain_prop.png", fontsize=30)
EQ["disagreement"] = render_eq(
    r"$\delta(t)=\sqrt{\frac{1}{N}\sum_i \|p_i(t)-\bar p(t)\|^2}$",
    "eq_disagreement.png", fontsize=22)
EQ["lambda2_zero"] = render_eq(
    r"$\lambda_2(L) = 0 \;\Rightarrow\; \mathrm{no\ global\ rendezvous}$",
    "eq_lambda2_zero.png", fontsize=24)
EQ["ratio"] = render_eq(
    r"$\tau/\tau_{\mathrm{crit}}$", "eq_ratio.png", fontsize=22)


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
def slide_01():
    s = new_slide()
    # accent side bar
    add_rect(s, 0, 0, Inches(0.16), SLIDE_H, fill=ACCENT_BLUE, line=None)
    add_text(s, Inches(1.0), Inches(2.15), Inches(11.0), Inches(0.4),
              "CONTROL OF MULTI-ROBOT SYSTEMS — PROJECT 2026",
              size=14, color=ACCENT_ORANGE, bold=True)
    add_text(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.0),
              "Delay-Aware Rendezvous of\nNonholonomic Multi-Robot Systems",
              size=38, color=WHITE, bold=True, line_spacing=1.05)
    add_rect(s, Inches(1.02), Inches(4.15), Inches(1.3), Pt(3.5), fill=ACCENT_BLUE, line=None)
    add_text(s, Inches(1.0), Inches(4.35), Inches(10.6), Inches(0.9),
              "Communication delays, nonholonomic constraints, and the role\n"
              "of the critical delay τcrit",
              size=18, color=MUTED, italic=True, line_spacing=1.2)
    add_text(s, Inches(1.0), Inches(6.2), Inches(10.6), Inches(0.9),
              "Samuele Civale   ·   Matteo Zamponi   ·   Roberto Passante",
              size=17, color=WHITE, bold=False)
    add_text(s, Inches(1.0), Inches(6.6), Inches(10.6), Inches(0.4),
              "Control of Multi-Robot Systems — Project 2026",
              size=13, color=MUTED)
    NOTES[1] = (
        "Good morning / good afternoon. We're presenting our project on delay-aware "
        "rendezvous of nonholonomic multi-robot systems. The central question we address "
        "is how much communication delay a distributed rendezvous controller can tolerate "
        "before the team loses collective stability, and how the theoretical consensus "
        "delay margin manifests itself once the agents are realistic nonholonomic unicycle "
        "robots rather than simple integrators. The work is joint between the three of us."
    )
    speaker_notes(s, NOTES[1])
    SLIDE_TITLES.append("Title")


# ============================================================================
# SLIDE 2 — PROBLEM & RESEARCH QUESTION
# ============================================================================
def slide_02():
    s = new_slide()
    y = add_title(s, "Problem & Research Question", kicker="Motivation")

    add_text(s, Inches(0.55), y, Inches(12.2), Inches(0.6),
              "A team of unicycle robots must rendezvous using only relative position "
              "information exchanged over a communication graph, delayed by τ.",
              size=17, color=MUTED, italic=True)

    # Main question banner
    qy = y + Inches(0.85)
    band = add_rect(s, Inches(0.55), qy, Inches(12.23), Inches(1.05), fill=PANEL2,
                     line=ACCENT_BLUE, line_w=1.5, rounded=True, radius=0.15)
    tf = band.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "How much delay can the distributed controller tolerate before rendezvous is lost?"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    # Three ingredients diagram
    dy = qy + Inches(1.45)
    box_w, box_h = Inches(3.15), Inches(1.55)
    gap = Inches(0.55)
    total_w = 3 * box_w + 2 * gap
    start_x = (SLIDE_W - total_w) // 2
    labels = [
        ("Graph topology", "L = D − A", ACCENT_BLUE),
        ("Communication delay", "τ (transport lag)", ACCENT_ORANGE),
        ("Nonholonomic dynamics", "unicycle (v, ω)", ACCENT_GREEN),
    ]
    centers = []
    for i, (label, sub, color) in enumerate(labels):
        x = start_x + i * (box_w + gap)
        bx = add_rect(s, x, dy, box_w, box_h, fill=PANEL, line=color, line_w=1.75,
                       rounded=True, radius=0.12)
        tf = bx.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(17)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r1.font.name = FONT
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(14)
        r2.font.color.rgb = color
        r2.font.name = FONT
        centers.append(x + box_w // 2)

    # plus signs between boxes
    for i in range(2):
        px = start_x + box_w + gap * i + (gap - Inches(0.3)) // 2 + i * box_w
        add_text(s, px, dy + box_h // 2 - Inches(0.2), Inches(0.3), Inches(0.4), "+",
                  size=26, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    # arrow down to conclusion pill
    ay = dy + box_h + Inches(0.12)
    add_arrow_down(s, (SLIDE_W - Inches(0.4)) // 2, ay, Inches(0.4), Inches(0.28), fill=MUTED)
    py = ay + Inches(0.34)
    add_pill(s, (SLIDE_W - Inches(4.6)) // 2, py, Inches(4.6), Inches(0.5),
              "Collective rendezvous stability", fill=PANEL2, border=ACCENT_GOLD,
              text_color=ACCENT_GOLD, size=15)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[2] = (
        "Rendezvous is the canonical multi-robot problem: a team must agree on a common "
        "meeting point using only local, relative measurements. We add two realistic "
        "ingredients. First, communication is delayed: sensing and transmission introduce a "
        "lag tau between when a neighbour's position is generated and when it's used. "
        "Second, the agents are unicycles, not integrators, so they cannot move "
        "instantaneously in an arbitrary direction. Our central question is how much delay "
        "the controller can tolerate before rendezvous breaks down, and it depends on three "
        "interacting ingredients: graph topology, delay, and nonholonomic dynamics. The talk "
        "develops the theory for a linear reference model, then tests it against the real "
        "nonlinear unicycle simulations."
    )
    speaker_notes(s, NOTES[2])
    SLIDE_TITLES.append("Problem & Research Question")


# ============================================================================
# SLIDE 3 — UNICYCLE + DISTRIBUTED CONSENSUS
# ============================================================================
def slide_03():
    s = new_slide()
    y = add_title(s, "Unicycle Model & Communication Graph", kicker="Preliminaries")

    left_w = Inches(6.3)
    add_text(s, Inches(0.55), y, left_w, Inches(0.3), "Unicycle kinematics", size=16,
              color=ACCENT_BLUE, bold=True)
    add_eq_picture(s, EQ["unicycle"], Inches(0.55), y + Inches(0.35), Inches(6.1), Inches(0.85),
                    align="left")

    add_text(s, Inches(0.55), y + Inches(1.35), left_w, Inches(0.3), "Graph Laplacian",
              size=16, color=ACCENT_BLUE, bold=True)
    add_eq_picture(s, EQ["laplacian"], Inches(0.55), y + Inches(1.7), Inches(2.2), Inches(0.55),
                    align="left")
    add_eq_picture(s, EQ["eig_order"], Inches(0.55), y + Inches(2.35), left_w, Inches(0.5),
                    align="left")

    ry = y + Inches(3.05)
    add_bullets(s, Inches(0.55), ry, left_w, Inches(2.2), [
        ("λ₂(L)  —  algebraic connectivity: graph is connected iff λ₂ > 0; "
         "sets the nominal convergence rate.", {"size": 15.5}),
        ("λmax(L)  —  governs the theoretical full-state delay margin (next slides).",
         {"size": 15.5}),
    ], bullet_color=ACCENT_BLUE, gap=0.22)

    # ring graph image on the right
    img_box_w, img_box_h = Inches(5.15), Inches(4.6)
    ix = Inches(7.6)
    iy = y - Inches(0.1)
    panel = add_rect(s, ix - Inches(0.15), iy - Inches(0.15), img_box_w + Inches(0.3),
                      img_box_h + Inches(0.55), fill=PANEL, line=LINE, rounded=True, radius=0.04)
    add_picture_contain(s, os.path.join(FIG, "u00_ring_graph.png"), ix, iy, img_box_w,
                         img_box_h - Inches(0.5))
    add_text(s, ix, iy + img_box_h - Inches(0.42), img_box_w, Inches(0.55),
              "Baseline graph: N = 6 ring   →   λ₂ = 1,  λmax = 4",
              size=14.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[3] = (
        "Each robot is a unicycle: it can only accelerate along its current heading and "
        "steer -- it cannot slide sideways. That's the nonholonomic constraint that will "
        "matter later. Communication is an undirected graph with Laplacian L equals D minus "
        "A, symmetric positive semi-definite, eigenvalues real and ordered from lambda_1 "
        "equals zero up to lambda_max. Two eigenvalues matter: lambda_2, the algebraic "
        "connectivity -- positive iff the graph is connected, sets the nominal convergence "
        "speed -- and lambda_max, which, as we derive in two slides, sets the theoretical "
        "delay margin. Our baseline is this six-node ring: lambda_2 equals 1, lambda_max "
        "equals 4."
    )
    speaker_notes(s, NOTES[3])
    SLIDE_TITLES.append("Unicycle Model & Communication Graph")


# ============================================================================
# SLIDE 4 — FROM CONSENSUS VECTOR TO A UNICYCLE
# ============================================================================
def slide_04():
    s = new_slide()
    y = add_title(s, "From the Consensus Field to the Unicycle", kicker="Modelling")

    add_text(s, Inches(0.55), y, Inches(12.2), Inches(0.55),
              "The delayed consensus protocol produces a desired planar velocity "
              "ui ∈ R². A unicycle cannot apply it directly: it can only move along its "
              "current heading.",
              size=16.5, color=MUTED, italic=True)

    # pipeline diagram
    py = y + Inches(0.85)
    box_h = Inches(1.05)
    w1, w2, w3 = Inches(2.5), Inches(3.6), Inches(2.6)
    gap = Inches(0.55)
    total = w1 + w2 + w3 + 2 * gap
    x0 = (SLIDE_W - total) // 2

    b1 = add_rect(s, x0, py, w1, box_h, fill=PANEL, line=ACCENT_ORANGE, line_w=1.75,
                   rounded=True, radius=0.14)
    tf = b1.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "desired velocity\nui ∈ R²"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT

    ax1 = x0 + w1 + Inches(0.06)
    add_arrow_right(s, ax1, py + box_h/2 - Inches(0.14), gap - Inches(0.12), Inches(0.28),
                     fill=MUTED)

    x2 = x0 + w1 + gap
    b2 = add_rect(s, x2, py, w2, box_h, fill=PANEL2, line=ACCENT_BLUE, line_w=1.75,
                   rounded=True, radius=0.14)
    tf = b2.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "force projection onto\nheading  ti = (cosθi, sinθi)"
    r.font.size = Pt(15.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT

    ax2 = x2 + w2 + Inches(0.06)
    add_arrow_right(s, ax2, py + box_h/2 - Inches(0.14), gap - Inches(0.12), Inches(0.28),
                     fill=MUTED)

    x3 = x2 + w2 + gap
    b3 = add_rect(s, x3, py, w3, box_h, fill=PANEL, line=ACCENT_GREEN, line_w=1.75,
                   rounded=True, radius=0.14)
    tf = b3.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "admissible unicycle\ninput (vi, ωi)"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT

    # equations below (default 'paper' force-projection controller)
    eqy = py + box_h + Inches(0.55)
    add_text(s, Inches(0.55), eqy, Inches(12.2), Inches(0.3),
              "Force-projection controller (default, CONTROLLER_TYPE = 'paper')",
              size=15, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_eq_picture(s, EQ["paper_def"], Inches(0.9), eqy + Inches(0.4), Inches(11.5), Inches(0.5))
    add_eq_picture(s, EQ["paper_law"], Inches(1.8), eqy + Inches(0.95), Inches(9.7), Inches(0.75))
    add_eq_picture(s, EQ["e_theta"], Inches(3.6), eqy + Inches(1.7), Inches(6.1), Inches(0.45))

    note = add_rect(s, Inches(0.55), eqy + Inches(2.3), Inches(12.23), Inches(0.55),
                     fill=PANEL2, line=ACCENT_GOLD, line_w=1.25, rounded=True, radius=0.2)
    tf = note.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left = Pt(14); tf.margin_right = Pt(14)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Only the projected, heading-aligned component of ui becomes thrust — this "
              "nonlinearity is why the linear consensus theorem does not transfer exactly.")
    r.font.size = Pt(14.5); r.font.italic = True; r.font.color.rgb = ACCENT_GOLD
    r.font.name = FONT

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[4] = (
        "Key modelling slide. The delayed consensus law generates, at every instant, a "
        "desired Cartesian velocity u_i in R^2. An integrator could apply it directly; a "
        "unicycle cannot -- it can only accelerate along its heading t_i. The default "
        "controller, the force-projection or 'paper' controller inspired by Listmann, "
        "Masalawala and Adamy, treats u_i as a force and projects it onto the heading: v_i "
        "is k_v times u_i dot t_i, so only the forward-aligned component produces thrust, "
        "while omega_i steers proportionally to the heading error e_theta. The orthogonal "
        "component is rejected and handled purely by steering. This projection is a genuine "
        "nonlinearity absent from the linear reference model -- the structural reason the "
        "critical delay we derive next cannot be an exact theorem for the nonlinear unicycle."
    )
    speaker_notes(s, NOTES[4])
    SLIDE_TITLES.append("From the Consensus Field to the Unicycle")


# ============================================================================
# SLIDE 5 — COMMUNICATION DELAY MODELS
# ============================================================================
def slide_05():
    s = new_slide()
    y = add_title(s, "Two Communication Delay Models", kicker="Delay modelling")

    col_w = Inches(5.85)
    gap = Inches(0.55)
    xL = Inches(0.55)
    xR = xL + col_w + gap
    col_h = Inches(4.35)

    # LEFT: full-state
    panelL = add_rect(s, xL, y, col_w, col_h, fill=PANEL, line=ACCENT_BLUE, line_w=1.75,
                       rounded=True, radius=0.06)
    add_text(s, xL + Inches(0.3), y + Inches(0.22), col_w - Inches(0.6), Inches(0.4),
              "FULL-STATE DELAY", size=17, color=ACCENT_BLUE, bold=True)
    add_eq_picture(s, EQ["fullstate"], xL + Inches(0.3), y + Inches(0.85),
                    col_w - Inches(0.6), Inches(0.7))
    add_bullets(s, xL + Inches(0.35), y + Inches(1.85), col_w - Inches(0.7), Inches(1.9), [
        ("The entire consensus state — including the agent's own position — "
         "enters the control law delayed.", {"size": 15.5}),
        ("No instantaneous self term: the only damping is on old information.",
         {"size": 15.5}),
    ], bullet_color=ACCENT_BLUE, gap=0.2)

    # RIGHT: neighbour-only
    panelR = add_rect(s, xR, y, col_w, col_h, fill=PANEL, line=ACCENT_ORANGE, line_w=1.75,
                       rounded=True, radius=0.06)
    add_text(s, xR + Inches(0.3), y + Inches(0.22), col_w - Inches(0.6), Inches(0.4),
              "NEIGHBOUR-ONLY DELAY", size=17, color=ACCENT_ORANGE, bold=True)
    add_eq_picture(s, EQ["neighboronly"], xR + Inches(0.25), y + Inches(0.8),
                    col_w - Inches(0.5), Inches(0.85))
    add_bullets(s, xR + Inches(0.35), y + Inches(1.85), col_w - Inches(0.7), Inches(1.9), [
        ("Own state pi(t) is current (locally available); only the communicated "
         "neighbour states pj are delayed.", {"size": 15.5}),
        ("Retains an instantaneous self-damping term ⇒ substantially more "
         "delay-tolerant, as the experiments confirm.", {"size": 15.5}),
    ], bullet_color=ACCENT_ORANGE, gap=0.2)

    add_eq_picture(s, EQ["neighboronly_stacked"], xR + Inches(0.25), y + Inches(3.55),
                    col_w - Inches(0.5), Inches(0.55))

    # bottom note
    ny = y + col_h + Inches(0.25)
    note = add_rect(s, Inches(0.55), ny, Inches(12.23), Inches(0.55), fill=PANEL2,
                     line=ACCENT_GOLD, line_w=1.25, rounded=True, radius=0.2)
    tf = note.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left = Pt(14); tf.margin_right = Pt(14)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "The theoretical τcrit derived next refers exactly to the full-state linear reference model."
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ACCENT_GOLD; r.font.name = FONT

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[5] = (
        "Before deriving the critical delay we must be precise about which delay model we "
        "mean -- the report studies two, and they behave very differently. Full-state delays "
        "everything, including the agent's own position: p dot equals minus k L p at t minus "
        "tau. Neighbour-only is more realistic: each robot knows its own current position "
        "instantly, so only communicated neighbour positions are delayed -- stacked form, p "
        "dot equals minus k times D p of t minus A p of t minus tau. The key difference is "
        "that undelayed term, minus k D p of t: an instantaneous damping the full-state model "
        "lacks, which is what makes neighbour-only far more delay-tolerant later. The "
        "critical delay we derive next is the exact threshold for full-state only; for "
        "neighbour-only it's a conservative reference, not an exact boundary."
    )
    speaker_notes(s, NOTES[5])
    SLIDE_TITLES.append("Two Communication Delay Models")


# ============================================================================
# SLIDE 6 — DERIVING THE CRITICAL DELAY
# ============================================================================
def slide_06():
    s = new_slide()
    y = add_title(s, "Deriving the Critical Delay τcrit", kicker="Main theory result")

    steps = [
        ("1 · Full-state protocol → modal decomposition",
         [EQ["step1a"], EQ["step1b"]], [Inches(3.1), Inches(4.3)]),
        ("2 · Try z(t) = e^{st}: characteristic equation", [EQ["step2"]], [Inches(7.6)]),
        ("3 · Stability boundary s = jω", [EQ["step3"]], [Inches(8.1)]),
        ("4 · Real / imaginary parts → first crossing", [EQ["step4"]], [Inches(7.8)]),
    ]

    rowy = y
    row_h = Inches(0.72)
    for label, imgs, widths in steps:
        add_text(s, Inches(0.55), rowy, Inches(4.3), row_h, label, size=13.5,
                  color=MUTED, italic=True, anchor=MSO_ANCHOR.MIDDLE)
        xcur = Inches(4.55)
        for img, w in zip(imgs, widths):
            add_eq_picture(s, img, xcur, rowy, w, row_h)
            xcur = xcur + w + Inches(0.15)
        rowy = rowy + row_h + Inches(0.06)

    add_eq_picture(s, EQ["step5"], Inches(0.55), rowy + Inches(0.05), Inches(12.2), Inches(0.55))

    # final boxed result
    fy = rowy + Inches(0.75)
    fw, fh = Inches(6.6), Inches(0.95)
    fx = (SLIDE_W - fw) // 2
    box = add_rect(s, fx, fy, fw, fh, fill=PANEL2, line=ACCENT_GOLD, line_w=2.25,
                    rounded=True, radius=0.18)
    add_eq_picture(s, EQ["taucrit_big"], fx + Inches(0.2), fy + Inches(0.08), fw - Inches(0.4),
                    fh - Inches(0.16))

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[6] = (
        "The full derivation, exact for the linear full-state model. Start from p dot equals "
        "minus k L p at t minus tau. Since L is real symmetric it has an orthonormal "
        "eigenbasis, so modal coordinates z equals V transpose p decouple the system into N "
        "scalar delay equations: z_i dot equals minus k lambda_i z_i at t minus tau. The zero "
        "eigenvalue mode is the conserved average; every other mode is z dot equals minus a z "
        "at t minus tau, a equals k lambda_i. Trying z of t equals e to the s t gives s "
        "equals minus a e to the minus s tau. Instability begins when a root first reaches "
        "the imaginary axis, s equals j omega; separating real and imaginary parts gives "
        "cosine omega tau equals zero, so omega tau equals pi over two at the first crossing, "
        "and omega equals a. Combining: tau_i equals pi over two k lambda_i per mode. The "
        "system is stable only while every mode is, so the binding constraint is the "
        "fastest mode, lambda_max -- giving tau critical equals pi over two k lambda_max of "
        "L. Physically, k lambda_max is the fastest rate the network imposes, and the margin "
        "is a quarter period of that mode's oscillation: a quarter-cycle delay turns "
        "corrective feedback into destabilizing, out-of-phase feedback."
    )
    speaker_notes(s, NOTES[6])
    SLIDE_TITLES.append("Deriving the Critical Delay")


# ============================================================================
# SLIDE 7 — BASELINE EXPERIMENT
# ============================================================================
def slide_07():
    s = new_slide()
    y = add_title(s, "Baseline Rendezvous (Sub-critical Delay)", kicker="Numerical validation")

    # left: video (trajectories, playable) + static disagreement plot
    img_w = Inches(3.75)
    img_h = Inches(3.35)
    gap = Inches(0.25)
    ix = Inches(0.55)
    add_rect(s, ix, y, img_w, img_h, fill=PANEL, line=LINE, rounded=True, radius=0.04)
    add_video(s, os.path.join(VID, "u01_full_state.mp4"),
              os.path.join(FIG, "u01_full_state_trajectories.png"), ix, y, img_w, img_h,
              label="Trajectories (animated .mp4 shipped alongside the deck)",
              frame_color=ACCENT_BLUE)
    ix2 = ix + img_w + gap
    add_picture_contain(s, os.path.join(FIG, "u01_full_state_disagreement.png"), ix2, y,
                         img_w, img_h)
    add_text(s, ix2, y + img_h, img_w, Inches(0.3), "RMS disagreement δ(t)", size=13,
              color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    # right: stats panel
    px = ix2 + img_w + Inches(0.35)
    pw = Inches(12.78) - px
    panel = add_rect(s, px, y, pw, img_h + Inches(0.3), fill=PANEL, line=LINE, rounded=True,
                      radius=0.06)
    add_text(s, px + Inches(0.25), y + Inches(0.18), pw - Inches(0.5), Inches(0.35),
              "Ring, N=6  ·  k=1", size=14, color=ACCENT_BLUE, bold=True)
    add_eq_picture(s, EQ["taucrit_ring"], px + Inches(0.2), y + Inches(0.55), pw - Inches(0.4),
                    Inches(0.5))
    stats = [
        ("τ used", "0.3 τcrit  ≈  0.118 s"),
        ("Final disagreement δ∞", "≈ 2.7 × 10⁻⁸ m"),
        ("Convergence time", "11.72 s"),
        ("Final centroid", "(0.356, 0.328)"),
    ]
    sy = y + Inches(1.2)
    for label, val in stats:
        add_text(s, px + Inches(0.25), sy, pw - Inches(0.5), Inches(0.3), label, size=12.5,
                  color=MUTED)
        add_text(s, px + Inches(0.25), sy + Inches(0.27), pw - Inches(0.5), Inches(0.35), val,
                  size=16.5, color=WHITE, bold=True)
        sy = sy + Inches(0.66)

    # bottom takeaway
    by = y + img_h + Inches(0.45)
    add_bullets(s, Inches(0.55), by, Inches(12.2), Inches(1.1), [
        ("With τ well below τcrit, the six unicycles rendezvous cleanly; disagreement decays "
         "monotonically to numerical zero.", {"size": 16}),
        ("This establishes the nominal, well-behaved regime used as a reference for every "
         "later experiment.", {"size": 16}),
    ], bullet_color=ACCENT_GREEN, gap=0.18)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[7] = (
        "Before stressing the system, we confirm nominal behaviour. On the ring, lambda_2 "
        "equals 1 and lambda_max equals 4, so with k equals 1 the reference critical delay "
        "is pi over eight, about 0.39 seconds. We run at 30 percent of that, tau equals "
        "0.118 seconds. On the left, six trajectories converge from scattered starting poses "
        "onto a common point, and RMS disagreement decays smoothly to about 8.3 times ten to "
        "the minus nine metres -- numerical zero. Convergence time is 6.22 seconds, settling "
        "near the initial centroid, (0.28, 0.11), consistent with undelayed theory. This is "
        "the nominal, healthy regime everything else in the talk pushes away from."
    )
    speaker_notes(s, NOTES[7])
    SLIDE_TITLES.append("Baseline Rendezvous")


# ============================================================================
# SLIDE 8 — THE CENTRAL EXPERIMENT: DELAY SWEEP
# ============================================================================
def slide_08():
    s = new_slide()
    y = add_title(s, "The Central Experiment: Delay Sweep", kicker="τ from 0 to 2 τcrit")

    img_w = Inches(5.95)
    img_h = Inches(3.05)
    ix = Inches(0.55)
    add_picture_contain(s, os.path.join(FIG, "u03_tau_sweep_summary.png"), ix, y, img_w, img_h)
    add_text(s, ix, y + img_h, img_w, Inches(0.25),
              "Final disagreement vs τ/τcrit (dashed line: τ/τcrit = 1)", size=12,
              color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    ix2 = ix + img_w + Inches(0.3)
    img_w2 = Inches(5.95)
    add_picture_contain(s, os.path.join(FIG, "u03_supracritical_curves.png"), ix2, y, img_w2,
                         img_h)
    add_text(s, ix2, y + img_h, img_w2, Inches(0.25), "Disagreement curves at τ = 1.2 τcrit",
              size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    by = y + img_h + Inches(0.32)
    add_bullets(s, Inches(0.55), by, Inches(12.2), Inches(2.0), [
        ("Full-state collapses near the predicted threshold: converges for τ ≤ 0.5 τcrit, "
         "fails from τ = 0.75 τcrit onward (δ∞: 0.61 → 1.08 m).",
         {"size": 14.5, "color": ACCENT_BLUE}),
        ("Neighbour-only is far more robust: still converges cleanly through 1.2 τcrit, "
         "degrading only near 2 τcrit.", {"size": 14.5, "color": ACCENT_ORANGE}),
        ("No exact transition at τ/τcrit = 1 for the nonlinear robot: the practical "
         "full-state crossing occurs somewhat earlier (projection nonlinearity, finite dt).",
         {"size": 14.5, "color": ACCENT_GOLD}),
    ], gap=0.12)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[8] = (
        "The single most important experiment. We sweep tau from zero to twice tau_crit for "
        "both delay models and plot final RMS disagreement against tau over tau_crit -- the "
        "dashed line marks the threshold. The blue full-state curve is near zero up to about "
        "half tau_crit, then jumps by orders of magnitude: 0.61 metres at 0.75 tau_crit, "
        "growing to 1.08 at twice tau_crit -- clear loss of rendezvous. Orange neighbour-only "
        "stays near numerical zero far beyond the threshold, still converging cleanly at 1.2 "
        "tau_crit, degrading only near 2 tau_crit. On the right, at 1.2 tau_crit: "
        "neighbour-only decays smoothly while full-state settles into a persistent "
        "oscillatory residual around 0.78 metres -- bounded, not diverging to infinity. One "
        "caution: the full-state practical crossing happens somewhat below the ideal ratio "
        "of one, between 0.5 and 0.75 tau_crit. That's expected -- tau_crit is exact only for "
        "the linear model, and the unicycle's heading dynamics, saturation and finite "
        "timestep erode that margin. We do not claim the nonlinear system transitions "
        "exactly at ratio one."
    )
    speaker_notes(s, NOTES[8])
    SLIDE_TITLES.append("The Central Experiment: Delay Sweep")


# ============================================================================
# SLIDE 9 — WHY TOPOLOGY MATTERS
# ============================================================================
def slide_09():
    s = new_slide()
    y = add_title(s, "Why Topology Matters", kicker="Graph spectrum & delay margin")

    img_w, img_h = Inches(6.7), Inches(4.35)
    add_picture_contain(s, os.path.join(FIG, "u04_topology_summary.png"), Inches(0.55), y,
                         img_w, img_h)
    add_text(s, Inches(0.55), y + img_h, img_w, Inches(0.3),
              "Convergence time by topology at τ = 0.5 τcrit (each graph normalised by its own τcrit)",
              size=12.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    px = Inches(7.6)
    pw = Inches(12.78) - px
    add_eq_picture(s, EQ["topology_prop"], px, y + Inches(0.1), pw, Inches(0.8))

    table_y = y + Inches(1.1)
    rows = [
        ("path", "3.73", "0.421"),
        ("ring", "4.00", "0.393"),
        ("star", "6.00", "0.262"),
        ("complete", "6.00", "0.262"),
    ]
    add_text(s, px, table_y, Inches(1.7), Inches(0.3), "graph", size=13, color=MUTED, bold=True)
    add_text(s, px + Inches(1.7), table_y, Inches(1.5), Inches(0.3), "λmax", size=13,
              color=MUTED, bold=True)
    add_text(s, px + Inches(3.2), table_y, Inches(1.9), Inches(0.3), "τcrit [s]", size=13,
              color=MUTED, bold=True)
    ry = table_y + Inches(0.35)
    for name, lmax, tcrit in rows:
        add_text(s, px, ry, Inches(1.7), Inches(0.32), name, size=15, color=WHITE, bold=True)
        add_text(s, px + Inches(1.7), ry, Inches(1.5), Inches(0.32), lmax, size=15, color=ACCENT_BLUE)
        add_text(s, px + Inches(3.2), ry, Inches(1.9), Inches(0.32), tcrit, size=15, color=ACCENT_ORANGE)
        ry = ry + Inches(0.4)

    by = ry + Inches(0.25)
    add_bullets(s, px, by, pw, Inches(1.6), [
        ("Denser graphs (star, complete) → larger λmax → smaller absolute delay margin.",
         {"size": 14.5}),
        ("Sparse path → smallest λmax → largest margin.", {"size": 14.5}),
        ("At equal τ/τcrit all connected topologies rendezvous; convergence speed orders "
         "itself by λ₂ (denser is faster).", {"size": 14.5}),
    ], gap=0.14, bullet_color=ACCENT_GOLD)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[9] = (
        "Tau_crit proportional to one over lambda_max tells us the delay margin is set by "
        "how dense the graph is. We compare path, ring, star and complete topologies, each "
        "run at half its own tau_crit for a fair relative comparison. The sparse path has "
        "the smallest lambda_max, 3.73, so the largest absolute margin, about 0.42 seconds; "
        "star and complete both reach lambda_max equals 6, the smallest margin, about 0.26 "
        "seconds. In the bar chart, every connected topology still rendezvouses at equal "
        "relative delay, but convergence times order themselves by lambda_2: denser graphs "
        "converge faster. So there's a real tension: a graph that produces fast, aggressive "
        "consensus is exactly the one that can least afford absolute delay. Lambda_2 buys "
        "you speed; lambda_max taxes your delay tolerance."
    )
    speaker_notes(s, NOTES[9])
    SLIDE_TITLES.append("Why Topology Matters")


# ============================================================================
# SLIDE 10 — CONTROLLER GAIN VS DELAY ROBUSTNESS
# ============================================================================
def slide_10():
    s = new_slide()
    y = add_title(s, "Controller Gain vs. Delay Robustness", kicker="Speed / robustness trade-off")

    img_w, img_h = Inches(6.7), Inches(4.35)
    add_picture_contain(s, os.path.join(FIG, "u11_gain_vs_taucrit.png"), Inches(0.55), y,
                         img_w, img_h)
    add_text(s, Inches(0.55), y + img_h, img_w, Inches(0.3),
              "τcrit and final disagreement vs. gain k, each run at fixed relative delay 0.8 τcrit(k)",
              size=12.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    px = Inches(7.6)
    pw = Inches(12.78) - px
    add_eq_picture(s, EQ["gain_prop"], px, y + Inches(0.1), pw, Inches(0.8))

    table_y = y + Inches(1.15)
    rows = [("0.5", "0.785"), ("1.0", "0.393"), ("1.5", "0.262"), ("2.0", "0.196")]
    add_text(s, px, table_y, Inches(2.0), Inches(0.3), "gain k", size=13, color=MUTED, bold=True)
    add_text(s, px + Inches(2.3), table_y, Inches(2.0), Inches(0.3), "τcrit [s]", size=13,
              color=MUTED, bold=True)
    ry = table_y + Inches(0.35)
    for k_val, tcrit in rows:
        add_text(s, px, ry, Inches(2.0), Inches(0.32), k_val, size=16, color=WHITE, bold=True)
        add_text(s, px + Inches(2.3), ry, Inches(2.0), Inches(0.32), tcrit, size=16,
                  color=ACCENT_BLUE)
        ry = ry + Inches(0.4)

    by = ry + Inches(0.3)
    add_bullets(s, px, by, pw, Inches(1.1), [
        ("τcrit ∝ 1/k — halving k doubles the admissible delay.", {"size": 14.5}),
        ("Higher gain → faster consensus, but smaller admissible delay.", {"size": 14.5}),
    ], bullet_color=ACCENT_ORANGE, gap=0.16)

    take = add_rect(s, px, by + Inches(1.3), pw, Inches(0.7), fill=PANEL2, line=ACCENT_GOLD,
                     line_w=1.25, rounded=True, radius=0.15)
    tf = take.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_right = Pt(12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Speed and delay robustness cannot be increased independently."
    r.font.size = Pt(15); r.font.bold = True; r.font.italic = True
    r.font.color.rgb = ACCENT_GOLD; r.font.name = FONT

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[10] = (
        "Now we isolate the consensus gain k, topology fixed at the ring. For each k we "
        "recompute tau_crit of k and run at a fixed relative delay, 0.8 times tau_crit of k. "
        "Doubling k from 0.5 to 1 halves tau_crit from 0.785 to 0.393 seconds, and doubling "
        "again to k equals 2 brings it to 0.196 -- exact inverse proportionality, confirming "
        "the formula. This is the practical statement of the speed-robustness trade-off: you "
        "cannot make the consensus law arbitrarily aggressive and arbitrarily delay-tolerant "
        "at the same time. Gain and topology are the two design levers, and both act through "
        "the same k times lambda_max product."
    )
    speaker_notes(s, NOTES[10])
    SLIDE_TITLES.append("Controller Gain vs. Delay Robustness")


# ============================================================================
# SLIDE 11 — FAILURE / NEGATIVE CONTROL: DISCONNECTED GRAPH
# ============================================================================
def slide_11():
    s = new_slide()
    y = add_title(s, "Negative Control: Disconnected Graph", kicker="Role of connectivity")

    img_w, img_h = Inches(5.7), Inches(4.5)
    add_picture_contain(s, os.path.join(FIG, "u08_disconnected_full_state.png"), Inches(0.55),
                         y, img_w, img_h)

    px = Inches(0.55) + img_w + Inches(0.35)
    pw = Inches(12.78) - px
    add_eq_picture(s, EQ["lambda2_zero"], px, y + Inches(0.15), pw, Inches(0.6))

    add_bullets(s, px, y + Inches(1.0), pw, Inches(3.2), [
        ("Graph: three disjoint communicating pairs, no edges between pairs.",
         {"size": 15.5}),
        ("The Laplacian's zero eigenvalue has multiplicity 3 (one per component): "
         "by the standard definition λ₂ = 0, so global rendezvous is impossible.",
         {"size": 15.5}),
        ("Each component still reaches its own local agreement: three clusters "
         "form instead of one common meeting point.", {"size": 15.5}),
        ("Negative control: connectivity, not the control law, is what enables "
         "global agreement.", {"size": 15.5, "color": ACCENT_GOLD, "bold": True}),
    ], bullet_color=ACCENT_RED, gap=0.2)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[11] = (
        "A deliberate negative control. We replace the ring with three disjoint pairs -- six "
        "robots, three two-node components, no edges between them. The Laplacian of a graph "
        "with three components has a zero eigenvalue of multiplicity three, so by the "
        "standard definition lambda_2 equals zero: not connected. The average is only "
        "conserved within each component, not across components, so global rendezvous is "
        "structurally impossible regardless of controller or delay. That's exactly what we "
        "see: each pair meets at its own point, three clusters instead of one. This confirms "
        "connectivity, captured by lambda_2, not the control law, is what enables global "
        "agreement at all."
    )
    speaker_notes(s, NOTES[11])
    SLIDE_TITLES.append("Negative Control: Disconnected Graph")


# ============================================================================
# SLIDE 12 — BEYOND BASIC RENDEZVOUS (2x2 grid)
# ============================================================================
def slide_12():
    s = new_slide()
    y = add_title(s, "Beyond Basic Rendezvous", kicker="Extensions of the delayed-consensus core")

    cell_w = Inches(5.9)
    cell_h = Inches(2.25)     # image area
    label_h = Inches(0.3)     # caption strip
    panel_h = cell_h + label_h
    gapx, gapy = Inches(0.35), Inches(0.12)
    row_pitch = panel_h + gapy
    x0 = Inches(0.55)
    items = [
        (os.path.join(FIG, "u15_leader_follower_trajectories.png"), "Leader–follower with moving target"),
        (os.path.join(FIG, "u16_formation_trajectories.png"), "Rigid formation (offset consensus)"),
        (os.path.join(FIG, "u19_avoidance_trajectories.png"), "Obstacle & inter-robot avoidance"),
        (os.path.join(FIG, "u20_switching_graph_trajectories.png"), "Switching communication graph"),
    ]
    for idx, (path, label) in enumerate(items):
        col = idx % 2
        row = idx // 2
        x = x0 + col * (cell_w + gapx)
        yy = y + row * row_pitch
        add_rect(s, x, yy, cell_w, panel_h, fill=PANEL, line=LINE, rounded=True, radius=0.05)
        add_picture_contain(s, path, x + Inches(0.08), yy + Inches(0.06), cell_w - Inches(0.16),
                             cell_h - Inches(0.1))
        add_text(s, x, yy + cell_h, cell_w, label_h, label, size=13.5,
                  color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[12] = (
        "The same delayed-consensus core extends beyond a single toy rendezvous. Top left: a "
        "leader tracks a moving circular reference while the rest follow through delayed "
        "neighbour-only consensus. Top right: consensus on shifted coordinates -- position "
        "minus a prescribed offset -- turns rendezvous into a rigid formation around an "
        "anchor point, the delay analysis carrying over unchanged. Bottom left: continuous "
        "repulsive potentials for an obstacle and for robot proximity let the team bend "
        "around danger while still rendezvousing with positive clearance. Bottom right: the "
        "graph is rebuilt every step from robot positions -- a switching proximity graph -- "
        "and the team still converges as long as it stays connected often enough. Together "
        "these show the core layer is not a fragile, single-purpose result."
    )
    speaker_notes(s, NOTES[12])
    SLIDE_TITLES.append("Beyond Basic Rendezvous")


# ============================================================================
# SLIDE 13 — ROBUSTNESS TO REALISTIC NETWORK EFFECTS
# ============================================================================
def slide_13():
    s = new_slide()
    y = add_title(s, "Robustness to Time-Varying Delay & Topology", kicker="Stress tests")

    img_h = Inches(2.95)
    w1 = Inches(3.85)
    add_picture_contain(s, os.path.join(FIG, "u22_crossing_critical_tau_profile.png"),
                         Inches(0.55), y, w1, img_h)
    add_text(s, Inches(0.55), y + img_h, w1, Inches(0.45),
              "τ(t)/τcrit: sinusoid peaking at 1.19, dipping to 0.71",
              size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    x2 = Inches(0.55) + w1 + Inches(0.25)
    w2 = Inches(3.85)
    add_picture_contain(s, os.path.join(FIG, "u22_crossing_critical_curves.png"), x2, y, w2, img_h)
    add_text(s, x2, y + img_h, w2, Inches(0.45),
              "Disagreement: full-state diverges, neighbour-only converges",
              size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    x3 = x2 + w2 + Inches(0.25)
    w3 = Inches(12.78) - x3
    add_picture_contain(s, os.path.join(FIG, "u20_switching_graph_lambda2.png"), x3, y, w3, img_h)
    add_text(s, x3, y + img_h, w3, Inches(0.45),
              "λ₂(L(t)) of the switching proximity graph — stays positive",
              size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    by = y + img_h + Inches(0.6)
    add_bullets(s, Inches(0.55), by, Inches(12.2), Inches(1.7), [
        ("Delay crossing τcrit: full-state final disagreement ≈0.74 m (no convergence); "
         "neighbour-only converges to ≈1×10⁻⁸ m in 9.5 s.",
         {"size": 15, "color": ACCENT_BLUE}),
        ("Neighbour-only tolerates excursions above the full-state reference threshold, "
         "consistent with its larger true margin.", {"size": 15, "color": ACCENT_ORANGE}),
        ("Switching topology: λ₂(t) fluctuates but stays positive, preserving convergence.",
         {"size": 15}),
    ], gap=0.13)

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[13] = (
        "Two stress tests probe the theory's assumptions. First, we make the delay itself "
        "time-varying: tau of t is a sinusoid around 0.95 tau_crit with 25 percent "
        "amplitude, swinging between 0.71 and 1.19 times tau_crit, repeatedly crossing the "
        "threshold -- left panel. Middle panel: full-state settles into a persistent "
        "disagreement around 0.74 metres and never converges, while neighbour-only still "
        "converges cleanly to numerical zero in about 9.5 seconds despite the excursions. "
        "Consistent with the earlier argument: neighbour-only's true margin, thanks to its "
        "self-damping, is larger than the full-state reference, so it absorbs crossings that "
        "break the full-state model. Second, on the right, the switching proximity graph's "
        "algebraic connectivity fluctuates but never touches zero, which is why the team "
        "still converges. Both reinforce the same message: neighbour-only is the more "
        "forgiving choice under realistic, non-ideal network conditions."
    )
    speaker_notes(s, NOTES[13])
    SLIDE_TITLES.append("Robustness to Time-Varying Delay & Topology")


# ============================================================================
# SLIDE 14 — CONCLUSIONS
# ============================================================================
def slide_14():
    s = new_slide()
    y = add_title(s, "Conclusions", kicker="Summary")

    add_eq_picture(s, EQ["taucrit_small"], Inches(8.4), Inches(0.5), Inches(4.3), Inches(0.55),
                    align="right")

    points = [
        ("1", "Communication delay fundamentally interacts with the graph spectrum",
         "the critical delay is governed jointly by the consensus gain k and by λmax(L).",
         ACCENT_BLUE),
        ("2", "Exact result for the linear full-state model",
         "τcrit = π / (2 k λmax(L)) is a proven stability threshold for the single-integrator "
         "delayed-consensus system.", ACCENT_ORANGE),
        ("3", "A reference, not a theorem, for the unicycle",
         "for nonholonomic unicycles τcrit remains a useful theoretical baseline, but it is "
         "not an exact nonlinear stability boundary — heading dynamics, projection and "
         "saturation shift the real transition.", ACCENT_GOLD),
        ("4", "Neighbour-only communication is markedly more robust",
         "its instantaneous self-state term supplies damping the full-state model lacks, so "
         "it tolerates delays well beyond the full-state reference threshold.", ACCENT_GREEN),
    ]

    cy = y + Inches(0.15)
    row_h = Inches(1.15)
    for num, head, body, color in points:
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), cy, Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background(); circ.shadow.inherit = False
        tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = BG; r.font.name = FONT

        add_text(s, Inches(1.35), cy - Inches(0.03), Inches(11.4), Inches(0.4), head, size=18,
                  color=WHITE, bold=True)
        add_text(s, Inches(1.35), cy + Inches(0.38), Inches(11.4), Inches(0.7), body, size=14.5,
                  color=MUTED, line_spacing=1.1)
        cy = cy + row_h

    take = add_rect(s, Inches(0.55), cy + Inches(0.05), Inches(12.23), Inches(0.55), fill=PANEL2,
                     line=ACCENT_GOLD, line_w=1.25, rounded=True, radius=0.2)
    tf = take.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left = Pt(14); tf.margin_right = Pt(14)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Fast consensus, dense connectivity and delay robustness must be designed together."
    r.font.size = Pt(15.5); r.font.bold = True; r.font.italic = True
    r.font.color.rgb = ACCENT_GOLD; r.font.name = FONT

    add_footer(s, "Delay-Aware Unicycle Rendezvous — Control of Multi-Robot Systems")
    NOTES[14] = (
        "To conclude. First, communication delay fundamentally interacts with the graph "
        "spectrum -- it's coupled to how the network is built and how aggressively it's "
        "driven. Second, for the linear full-state model we derived the exact threshold "
        "tau_crit equals pi over two k lambda_max of L, verified numerically to separate "
        "convergent from divergent regimes. Third, for nonholonomic unicycles tau_crit "
        "remains a useful theoretical reference -- behaviour changes character near it -- "
        "but it is not an exact nonlinear stability boundary, since the force-projection "
        "controller, heading dynamics and saturation are real nonlinearities absent from the "
        "linear model. Fourth, neighbour-only communication is considerably more robust in "
        "every experiment, since each robot's own current position is always locally "
        "available, supplying damping the full-state model lacks. The overarching message: "
        "fast consensus, dense connectivity and delay robustness must be designed together, "
        "and where the architecture allows it, neighbour-only is the safer choice. Thank you "
        "-- happy to take questions."
    )
    speaker_notes(s, NOTES[14])
    SLIDE_TITLES.append("Conclusions")


# ============================================================================
# BUILD
# ============================================================================
def build():
    slide_01(); slide_02(); slide_03(); slide_04(); slide_05(); slide_06()
    slide_07(); slide_08(); slide_09(); slide_10(); slide_11(); slide_12()
    slide_13(); slide_14()
    prs.save(OUT_PPTX)
    print(f"Saved: {OUT_PPTX}")
    print(f"Slides: {len(prs.slides)}")
    write_notes_md()
    write_outline_md()
    return NOTES, SLIDE_TITLES


def write_notes_md():
    out = os.path.join(HERE, "slide_notes.md")
    lines = ["# Speaker Notes — Delay-Aware Rendezvous of Nonholonomic Multi-Robot Systems\n",
             "Target: ~10-12 minutes total. Each block below is ~40-70 seconds of spoken material.\n"]
    for i in range(1, len(SLIDE_TITLES) + 1):
        lines.append(f"## Slide {i} — {SLIDE_TITLES[i-1]}\n")
        lines.append(NOTES.get(i, "").strip() + "\n")
    with open(out, "w") as f:
        f.write("\n".join(lines))
    print(f"Saved: {out}")


def write_outline_md():
    out = os.path.join(HERE, "slide_outline.md")
    lines = ["# Slide Outline — Delay-Aware Rendezvous of Nonholonomic Multi-Robot Systems\n",
             f"14 slides, 16:9, ~10-12 minute talk.\n"]
    for i, title in enumerate(SLIDE_TITLES, 1):
        lines.append(f"{i}. {title}")
    with open(out, "w") as f:
        f.write("\n".join(lines) + "\n")
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
